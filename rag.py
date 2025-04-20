#改用api版
import os

os.environ["TOKENIZERS_PARALLELISM"] = "false"  # 禁用 Hugging Face 的并行警告
os.environ["BITSANDBYTES_NOWELCOME"] = "1"  # 禁止 bitsandbytes 欢迎信息
os.environ["WANDB_DISABLED"] = "true"  # 禁用 wandb 日志（如果有）

import faiss
import numpy as np
import pickle
import torch
from transformers import AutoModelForCausalLM, AutoTokenizer
from sentence_transformers import SentenceTransformer
import pandas as pd
from config import config
import logging
from docx import Document  # 以防后面需要 docx 解析

from openai import OpenAI
# 初始化 API 客户端
client = OpenAI(api_key=config.API_KEY, base_url=config.BASE_URL)

##############################################################################
# 1. 加载模型与索引
##############################################################################

# 使用 API 获取嵌入向量，调用 config 中的 get_embedding 方法
def get_embedding_from_api(text):
    return np.array(config.get_embedding(text))  # 调用 config 中的方法来获取嵌入向量


# 加载 FAISS 索引
index = faiss.read_index(str(config.FAISS_CACHE / "docs.index"))
print(f"✅ FAISS 索引维度：{index.d}")

# 加载文件名列表
with open(str(config.FAISS_CACHE / "filenames.pkl"), "rb") as f:
    filenames = pickle.load(f)

# 获取 LLM 响应的功能，现在通过调用 API 来获取大模型的回答
def get_llm_response_from_api(prompt):
    return config.get_llm_response(prompt)  # 使用 config 中的方法来获取 LLM 响应

# 使用 API 获取嵌入和模型响应
sample_text = "这是一段用于测试的文本"
embedding = get_embedding_from_api(sample_text)
print(f"✅ 获取的嵌入向量：{embedding[:5]}...")  # 输出前 5 个元素作为示例

# 通过 API 获取 LLM 模型的响应
prompt = "请生成一段关于人工智能的介绍。"
llm_response = get_llm_response_from_api(prompt)
print(f"✅ LLM 模型的响应：{llm_response}")


##############################################################################
# 2. 文档加载函数
##############################################################################

def load_document_content(filename):
    """
    根据文件名加载文本内容。identity.md 作为 AI 角色设定。
    """
    # 使用配置中的知识库目录
    file_path = os.path.join(config.REFERENCE_FOLDER, filename)

    # 特殊处理 identity.md
    if filename == "identity.md":
        with open(file_path, "r", encoding="utf-8") as file:
            return f"[系统角色提示]\n{file.read()}\n\n"

    if filename.endswith(".txt") or filename.endswith(".md"):
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    elif filename.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif filename.endswith(".docx"):
        # 利用 python-docx 读取
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    elif filename.endswith(".xlsx") or filename.endswith(".csv"):
        try:
            df = pd.read_excel(file_path) if filename.endswith(".xlsx") else pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            return f"❌ 读取 {filename} 时出错: {str(e)}"
    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            return f"❌ 读取 {filename} 时出错: {str(e)}"

    return "❌ 无法读取该文档格式：" + filename


##############################################################################
# 3. 检索函数
##############################################################################

def retrieve_top_k_documents(query, k=3):
    """
    根据查询语句在索引中找到最相关的 k 个文档，并返回其内容（截取）。
    """
    query_embedding = get_embedding_from_api(query)  # 获取嵌入向量

    # 确保查询向量是二维数组 (1, d)
    query_embedding = np.expand_dims(query_embedding, axis=0)

    print(f"✅ 查询向量维度：{query_embedding.shape[0]}")

    _, idxs = index.search(query_embedding, k)

    retrieved_docs = []
    for i in idxs[0]:
        filename = filenames[i]
        content = load_document_content(filename)
        # 这里可以只截取前1000字符，避免 Prompt 过长
        retrieved_docs.append(f"📄【{filename}】\n{content[:100]}...")
    return retrieved_docs


##############################################################################
# 4. 分块 (Chunk) + 摘要 / 改写
##############################################################################

def chunk_text(text, chunk_size=1000, overlap=200):
    """
    将长文本分块，每块 chunk_size 个字符，并在块之间保留 overlap 个字符重叠，避免关键信息被切断。
    """
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)  # 避免越界
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap

        if start < 0:
            start = 0
        if start >= text_len:
            break
    return chunks


def summarize_long_text(text):
    """
    对长文本进行多段式摘要，然后合并。
    """
    # 分块
    text_chunks = chunk_text(text, chunk_size=1500, overlap=200)
    chunk_summaries = []

    # 逐块摘要
    for idx, chunk in enumerate(text_chunks):
        prompt = f"请阅读以下文本内容，并进行简要总结：\n{chunk}\n\n总结："
        llm_response = get_llm_response_from_api(prompt)  # 通过 API 获取响应
        chunk_summaries.append(llm_response.strip())

    # 合并所有块的摘要，再让模型做一次“总总结”
    combined_summary = "\n".join(chunk_summaries)
    final_prompt = f"以下是多个分块的总结，请将其合并为一个简洁的整体总结：\n{combined_summary}\n\n整体总结："
    final_summary = get_llm_response_from_api(final_prompt)  # 通过 API 获取总结
    return final_summary.strip()


def rewrite_long_text(text):
    """
    对长文本进行“改写 / 润色”，与摘要类似的思路，先分块再合并。
    """
    # 先粗暴示例，也可以根据需要拆分做多次改写
    text_chunks = chunk_text(text, chunk_size=1500, overlap=200)
    rewrite_results = []

    for chunk in text_chunks:
        prompt = f"请对以下文本进行语言润色或改写，使其更通顺、简洁：\n{chunk}\n\n改写后："
        rewrite_result = get_llm_response_from_api(prompt)  # 通过 API 获取改写
        rewrite_results.append(rewrite_result.strip())

    # 简单拼接，如果想再做最终合并，可以再来一次生成
    return "\n".join(rewrite_results).strip()




##############################################################################
# 5. 最终回答：generate_answer
#    - 如果用户问 "总结xx文件" 或 "改写xx文件"：
#         -> 直接找到文件内容做 summarize/rewrite
#    - 否则做普通RAG问答
##############################################################################

def generate_answer(query):
    # 1) 检索文档
    related_docs = retrieve_top_k_documents(query, k=3)

    # 2) 加载 identity.md（如果存在）
    identity_content = ""
    if "identity.md" in filenames:
        identity_content = load_document_content("identity.md")

    # 3) 拼接上下文
    context = f"AI Identity/Persona\n{identity_content}\n\n【知识库检索结果】\n" + "\n\n".join(related_docs)

    # 构造 Prompt
    prompt = f"""
你是一名智能问答助手(扮演DeepSeek知识管家)，以下是你的身份描述和检索到的文档内容:
{context}

请遵守identity.md中的所有“沟通风格准则”和“特殊行为准则”，并根据文档做出回答。
如果在知识库中找不到答案，请回答“对不起，我在知识库中没有找到相关信息”。
请使用简洁且专业的口吻。
用户的问题：{query}

请直接给出简洁且专业的回答：
""".strip()

    # 通过 API 获取流式回答
    llm_response = get_llm_response_from_api(prompt)
    return llm_response.strip()


##############################################################################
# 6. 简单函数：_simple_summarize / _simple_rewrite (对短文本)
##############################################################################

def _simple_summarize(text):
    """
    对短文本做一次性摘要。如果文本不大，可直接用。
    """
    prompt = f"请阅读以下文本并进行简要总结：\n{text}\n\n总结："
    llm_response = get_llm_response_from_api(prompt)  # 通过 API 获取摘要
    return llm_response.strip()


def _simple_rewrite(text):
    """
    对短文本做一次性改写。
    """
    prompt = f"请对以下文本进行语言润色或改写，使其更通顺、简洁：\n{text}\n\n改写后："
    llm_response = get_llm_response_from_api(prompt)  # 通过 API 获取改写
    return llm_response.strip()


##############################################################################
# 7. 命令行交互入口
##############################################################################


import re

if __name__ == "__main__":
    print("📚 RAG + Summarize/Rewrite 示例程序启动...")
    try:
        while True:
            query = input("\n请输入您的问题（输入 'exit' 退出）： ")
            if query.lower().strip() == "exit":
                print("\n👋 退出 RAG 系统，再见！")
                break

            # 初始化变量
            bot_response = ""
            debug_info = "暂无推理过程"
            current_output_buffer = ""

            # 流式处理生成器输出
            stream = generate_answer(query)
            for current_output in stream:
                current_output_buffer += current_output
                if "</think>" in current_output_buffer:
                    parts = current_output_buffer.split("</think>")
                    if len(parts) >= 2:
                        bot_response = parts[1].strip()  # 模型的回答
                        if len(parts) > 2:
                            debug_info = parts[2].strip()  # 调试信息
                        else:
                            debug_info = "暂无推理过程"
                        # 清空缓冲区
                        current_output_buffer = ""
                        # 打印流式输出的回答
                        print("\nAI 回答（流式输出）：", bot_response, "\n")
                    else:
                        bot_response = parts[0].strip()
                        current_output_buffer = ""
                else:
                    # 如果没有找到 </think>，继续积累输出
                    continue

            # 如果缓冲区中还有内容，处理剩余部分
            if current_output_buffer:
                bot_response = current_output_buffer.strip()
                print("\nAI 回答（流式输出）：", bot_response, "\n")

    except KeyboardInterrupt:
        print("\n👋 检测到 Ctrl + C，退出 RAG 系统！")

"""  模型回答 一次性输出
if __name__ == "__main__":
    print("📚 RAG + Summarize/Rewrite 示例程序启动...")
    try:
        while True:
            query = input("\n请输入您的问题（输入 'exit' 退出）： ")
            if query.lower().strip() == "exit":
                print("\n👋 退出 RAG 系统，再见！")
                break

            # 初始化变量
            bot_response = ""
            debug_info = "暂无推理过程"

            # 流式处理生成器输出
            stream = generate_answer(query)
            for current_output in stream:
                if "</think>" in current_output:
                    parts = current_output.split("</think>")
                    if len(parts) >= 2:
                        bot_response = parts[1].strip()  # 模型的回答
                        if len(parts) > 2:
                            debug_info = parts[2].strip()  # 调试信息
                    else:
                        bot_response = parts[0].strip()
                else:
                    bot_response = current_output.strip()

            # 输出最终回答
            print("\nAI 回答：", bot_response, "\n")

    except KeyboardInterrupt:
        print("\n👋 检测到 Ctrl + C，退出 RAG 系统！")
"""

